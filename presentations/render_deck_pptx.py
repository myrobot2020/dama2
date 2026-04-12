"""
Parse presentations/DAMA_deck.marp.md (Marp-style) and write DAMA_deck.pptx.
Run: python presentations/render_deck_pptx.py
Requires: pip install python-pptx
"""
from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
MD_PATH = ROOT / "DAMA_deck.marp.md"
OUT_PATH = ROOT / "DAMA_deck.pptx"

HEADER_FOOTER = "DAMA · April 2026  ·  Dhamma Alignment & Verification Architecture"


def strip_comments(s: str) -> str:
    return re.sub(r"<!--.*?-->", "", s, flags=re.DOTALL).strip()


def split_front_matter(text: str) -> str:
    lines = text.splitlines()
    if not lines or lines[0].strip() != "---":
        return text
    i = 1
    while i < len(lines) and lines[i].strip() != "---":
        i += 1
    if i >= len(lines):
        return text
    return "\n".join(lines[i + 1 :]).lstrip("\n")


def split_slides(body: str) -> list[str]:
    parts = re.split(r"\n---\n", body)
    return [strip_comments(p).strip() for p in parts if strip_comments(p).strip()]


def parse_inline_runs(paragraph, line: str) -> None:
    """Split on **bold** and apply font.bold to those segments."""
    paragraph.clear()
    parts = re.split(r"(\*\*.+?\*\*|\*.+?\*)", line)
    for part in parts:
        if not part:
            continue
        run = paragraph.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]
            run.font.bold = True
        elif part.startswith("*") and part.endswith("*") and len(part) > 2:
            run.text = part[1:-1]
            run.font.italic = True
        else:
            run.text = part


def slide_title_and_bullets(prs: Presentation, raw: str) -> None:
    lines = [ln.rstrip() for ln in raw.splitlines()]
    while lines and not lines[0].strip():
        lines.pop(0)
    if not lines:
        return

    title_line = lines[0].strip()
    if title_line.startswith("# "):
        title = title_line[2:].strip()
        body_lines = lines[1:]
    elif title_line.startswith("## "):
        title = title_line[3:].strip()
        body_lines = lines[1:]
    else:
        title = "Slide"
        body_lines = lines

    layout = prs.slide_layouts[1]  # title and content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    tf = slide.placeholders[1].text_frame
    tf.clear()
    first = True
    for line in body_lines:
        line_stripped = line.strip()
        if not line_stripped:
            continue
        if line_stripped.startswith("- "):
            text = line_stripped[2:].strip()
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.level = 0
            parse_inline_runs(p, text)
        elif re.match(r"^\d+\.\s", line_stripped):
            text = re.sub(r"^\d+\.\s*", "", line_stripped)
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.level = 0
            parse_inline_runs(p, text)
        else:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.level = 0
            parse_inline_runs(p, line_stripped)

    _add_footer_note(slide, prs)


def slide_title_only_centered(prs: Presentation, raw: str) -> None:
    """First slide style: # title + body lines in one text box."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    left, top, width, height = Inches(0.6), Inches(1.0), Inches(9.0), Inches(5.5)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True

    lines = [ln.strip() for ln in raw.splitlines() if ln.strip()]
    if not lines:
        return

    if lines[0].startswith("## "):
        title = lines[0][3:].strip()
        rest = lines[1:]
    elif lines[0].startswith("# "):
        title = lines[0][2:].strip()
        rest = lines[1:]
    elif lines[0].startswith("**") and lines[0].endswith("**"):
        title = lines[0][2:-2].strip()
        rest = lines[1:]
    else:
        title = lines[0].lstrip("#").strip()
        rest = lines[1:]

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    p.alignment = PP_ALIGN.CENTER

    for line in rest:
        p = tf.add_paragraph()
        if line.startswith("**") and line.endswith("**"):
            p.text = line[2:-2]
            p.font.bold = True
        else:
            parse_inline_runs(p, line)
        p.font.size = Pt(22)
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(8)

    _add_footer_note(slide, prs)


def _add_footer_note(slide, prs: Presentation) -> None:
    h = prs.slide_height
    box = slide.shapes.add_textbox(Inches(0.5), h - Inches(0.55), Inches(9.0), Inches(0.45))
    p = box.text_frame.paragraphs[0]
    p.text = HEADER_FOOTER
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)


def main() -> int:
    if not MD_PATH.is_file():
        print(f"Missing {MD_PATH}", file=sys.stderr)
        return 1

    text = MD_PATH.read_text(encoding="utf-8")
    body = split_front_matter(text)
    slides_raw = split_slides(body)
    if not slides_raw:
        print("No slides found.", file=sys.stderr)
        return 1

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for raw in slides_raw:
        # Section slides: first line is `## Title`. Lead / closing slides use `#` or `**…**`.
        if raw.lstrip().startswith("## "):
            slide_title_and_bullets(prs, raw)
        else:
            slide_title_only_centered(prs, raw)

    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
